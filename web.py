"""
skill-forge/web.py
网页版 Skill Forge（支持多轮对话 + 文件上传 + 文件下载）
"""

import os
import io
import json
import base64
import re
from datetime import datetime
from dotenv import load_dotenv
from openai import OpenAI
import streamlit as st

# ============ 1. 初始化 ============

load_dotenv()

client = OpenAI(
    api_key=os.getenv("DEEPSEEK_API_KEY"),
    base_url="https://api.deepseek.com"
)

SKILLS_DIR = "skills"
OUTPUT_DIR = "outputs"
if not os.path.exists(SKILLS_DIR):
    os.makedirs(SKILLS_DIR)
if not os.path.exists(OUTPUT_DIR):
    os.makedirs(OUTPUT_DIR)

if "sop" not in st.session_state:
    st.session_state.sop = None
if "sop_history" not in st.session_state:
    st.session_state.sop_history = []
if "skill" not in st.session_state:
    st.session_state.skill = None
if "chat_history" not in st.session_state:
    st.session_state.chat_history = []
if "uploaded_text" not in st.session_state:
    st.session_state.uploaded_text = ""


# ============ 2. 文件读取 ============

def read_uploaded_file(uploaded_file):
    """读取上传的文件，提取文本内容"""
    name = uploaded_file.name.lower()

    try:
        if name.endswith(".txt") or name.endswith(".md"):
            return uploaded_file.read().decode("utf-8")

        elif name.endswith(".json"):
            data = json.loads(uploaded_file.read().decode("utf-8"))
            return json.dumps(data, ensure_ascii=False, indent=2)

        elif name.endswith(".csv"):
            return uploaded_file.read().decode("utf-8")

        elif name.endswith(".docx"):
            from docx import Document
            doc = Document(uploaded_file)
            return "\n".join([p.text for p in doc.paragraphs if p.text.strip()])

        elif name.endswith(".xlsx") or name.endswith(".xls"):
            import openpyxl
            wb = openpyxl.load_workbook(uploaded_file)
            text = ""
            for sheet in wb.sheetnames:
                ws = wb[sheet]
                text += f"\n--- 工作表: {sheet} ---\n"
                for row in ws.iter_rows(values_only=True):
                    cells = [str(c) if c is not None else "" for c in row]
                    text += " | ".join(cells) + "\n"
            return text

        elif name.endswith(".pptx"):
            from pptx import Presentation
            prs = Presentation(uploaded_file)
            text = ""
            for i, slide in enumerate(prs.slides):
                text += f"\n--- 幻灯片 {i+1} ---\n"
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += shape.text + "\n"
            return text

        elif name.endswith(".pdf"):
            try:
                import PyPDF2
                reader = PyPDF2.PdfReader(uploaded_file)
                text = ""
                for page in reader.pages:
                    text += page.extract_text() + "\n"
                return text
            except ImportError:
                return "[PDF 文件需要安装 PyPDF2：pip install PyPDF2]"

        else:
            try:
                return uploaded_file.read().decode("utf-8")
            except Exception:
                return f"[无法读取 {name}，不支持该格式]"

    except Exception as e:
        return f"[读取 {name} 失败：{e}]"


# ============ 3. 文件生成 ============

def generate_txt(content, filename):
    """生成 TXT 文件"""
    return content.encode("utf-8"), f"{filename}.txt", "text/plain"


def generate_word(content, filename):
    """生成 Word 文件"""
    from docx import Document
    doc = Document()
    for line in content.split("\n"):
        line = line.strip()
        if not line:
            continue
        if line.startswith("# "):
            doc.add_heading(line[2:], level=1)
        elif line.startswith("## "):
            doc.add_heading(line[3:], level=2)
        elif line.startswith("### "):
            doc.add_heading(line[4:], level=3)
        else:
            doc.add_paragraph(line)
    buf = io.BytesIO()
    doc.save(buf)
    buf.seek(0)
    return buf.getvalue(), f"{filename}.docx", "application/vnd.openxmlformats-officedocument.wordprocessingml.document"


def generate_excel(content, filename):
    """生成 Excel 文件，content 应该是表格文本"""
    import openpyxl
    wb = openpyxl.Workbook()
    ws = wb.active
    ws.title = "Sheet1"
    for row_idx, line in enumerate(content.strip().split("\n"), 1):
        line = line.strip()
        if not line:
            continue
        line = line.strip("|")
        cells = [c.strip() for c in line.split("|")]
        for col_idx, cell in enumerate(cells, 1):
            if cell.replace("-", "").strip() == "":
                continue
            ws.cell(row=row_idx, column=col_idx, value=cell)
    buf = io.BytesIO()
    wb.save(buf)
    buf.seek(0)
    return buf.getvalue(), f"{filename}.xlsx", "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet"


def generate_ppt(content, filename):
    """生成 PPT 文件"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()

    slides_text = content.split("---")
    if len(slides_text) == 1:
        slides_text = content.split("\n\n")

    for slide_text in slides_text:
        slide_text = slide_text.strip()
        if not slide_text:
            continue
        slide = prs.slides.add_slide(prs.slide_layouts[1])
        lines = slide_text.split("\n")
        title_text = lines[0].lstrip("#").strip() if lines else "幻灯片"
        body_text = "\n".join(lines[1:]).strip() if len(lines) > 1 else ""
        slide.shapes.title.text = title_text
        if body_text and slide.placeholders[1]:
            slide.placeholders[1].text = body_text

    buf = io.BytesIO()
    prs.save(buf)
    buf.seek(0)
    return buf.getvalue(), f"{filename}.pptx", "application/vnd.openxmlformats-officedocument.presentationml.presentation"


def auto_generate_file(content, output_format, skill_name):
    """根据指定格式自动生成文件"""
    safe_name = skill_name.replace(" ", "_").replace("/", "_")
    timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    filename = f"{safe_name}_{timestamp}"

    fmt = output_format.lower().strip()

    if fmt in ["word", "docx", ".docx"]:
        return generate_word(content, filename)
    elif fmt in ["excel", "xlsx", ".xlsx"]:
        return generate_excel(content, filename)
    elif fmt in ["ppt", "pptx", ".pptx"]:
        return generate_ppt(content, filename)
    elif fmt in ["txt", ".txt", "text"]:
        return generate_txt(content, filename)
    elif fmt in ["json", ".json"]:
        return content.encode("utf-8"), f"{filename}.json", "application/json"
    elif fmt in ["md", "markdown", ".md"]:
        return content.encode("utf-8"), f"{filename}.md", "text/markdown"
    else:
        return generate_txt(content, filename)


# ============ 4. SOP 生成 ============

def call_generate_sop(task_description, deliverable):
    prompt = f"""你是一个专业的流程设计专家。请根据以下信息，生成一份详细的标准操作流程(SOP)。

## 任务描述
{task_description}

## 交付要求
{deliverable}

请以 JSON 格式输出，结构如下：
{{
    "title": "SOP标题",
    "objective": "目标概述",
    "steps": [
        {{
            "step_number": 1,
            "title": "步骤标题",
            "description": "具体做什么、怎么做",
            "input": "这一步需要什么",
            "output": "这一步产出什么",
            "acceptance_criteria": "怎么算做完了"
        }}
    ],
    "quality_checklist": ["检查项1", "检查项2"],
    "final_deliverable": "最终交付物描述"
}}

要求：
1. 步骤要细致，每一步都是可执行的
2. 上一步的 output 要能衔接下一步的 input
3. 每步都有明确的完成标准

只输出JSON，不要其他内容。"""

    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        response_format={"type": "json_object"}
    )
    return json.loads(response.choices[0].message.content)


# ============ 5. SOP 修改 ============

def call_refine_sop(current_sop, feedback):
    prompt = f"""你之前生成了以下 SOP：

{json.dumps(current_sop, ensure_ascii=False, indent=2)}

用户的反馈是：
{feedback}

请根据反馈修改 SOP，输出修改后的完整 JSON（格式不变）。
只输出 JSON，不要其他内容。"""

    response = client.chat.completions.create(
        model="deepseek-chat",
        messages=[{"role": "user", "content": prompt}],
        temperature=0.3,
        response_format={"type": "json_object"}
    )
    return json.loads(response.choices[0].message.content)


# ============ 6. 生成 Skill ============

def call_generate_skill(sop):
    prompt_for_system = f"""请根据以下 SOP，为一个 AI 助手编写 system prompt。
这个 AI 助手未来会按照这个 SOP 自动执行任务。

SOP 内容：
{json.dumps(sop, ensure_ascii=False, indent=2)}

要求：
1. system prompt 要包含完整的执行流程
2. 要包含每一步的具体操作指引
3. 要包含质量检查环节
4. 要告诉 AI 以什么格式输出结果
5. 要专业、清晰、无歧义

直接输出 system prompt 文本，不要任何包装。"""

    r1 = client.chat.completions.create(
        model="deepseek-chat",
        messages=[{"role": "user", "content": prompt_for_system}],
        temperature=0.2
    )
    system_prompt = r1.choices[0].message.content

    prompt_for_schema = f"""根据以下 SOP，定义这个工具的输入参数和输出格式。

SOP 内容：
{json.dumps(sop, ensure_ascii=False, indent=2)}

请以 JSON 格式输出：
{{
    "input_params": [
        {{
            "name": "参数名",
            "description": "参数描述",
            "type": "string",
            "required": true,
            "example": "示例值"
        }}
    ],
    "output_format": {{
        "description": "输出描述",
        "fields": [
            {{"name": "字段名", "description": "字段描述"}}
        ]
    }}
}}

只输出 JSON。"""

    r2 = client.chat.completions.create(
        model="deepseek-chat",
        messages=[{"role": "user", "content": prompt_for_schema}],
        temperature=0.2,
        response_format={"type": "json_object"}
    )
    schema = json.loads(r2.choices[0].message.content)

    skill = {
        "skill_name": sop["title"],
        "description": sop["objective"],
        "version": "1.0",
        "created_at": datetime.now().strftime("%Y-%m-%d %H:%M:%S"),
        "system_prompt": system_prompt,
        "input_params": schema.get("input_params", []),
        "output_format": schema.get("output_format", {}),
        "source_sop": sop
    }
    return skill


# ============ 7. 显示 SOP ============

def display_sop(sop):
    st.markdown(f"## 📋 {sop['title']}")
    st.markdown(f"**🎯 目标：** {sop['objective']}")
    st.markdown("---")
    for step in sop["steps"]:
        st.markdown(f"### 步骤 {step['step_number']}：{step['title']}")
        st.markdown(f"- 📖 **描述：** {step['description']}")
        st.markdown(f"- 📥 **输入：** {step['input']}")
        st.markdown(f"- 📤 **输出：** {step['output']}")
        st.markdown(f"- ✅ **完成标准：** {step['acceptance_criteria']}")
        st.markdown("")
    st.markdown("---")
    st.markdown("### 🔍 质量检查清单")
    for item in sop["quality_checklist"]:
        st.markdown(f"- {item}")
    st.markdown(f"**📦 最终交付物：** {sop['final_deliverable']}")


# ============ 8. 网页界面 ============

st.set_page_config(page_title="Skill Forge", page_icon="🔧", layout="wide")

st.title("🔧 Skill Forge")
st.markdown("*输入任务描述 → AI 生成 SOP → 你确认修改 → 固化为可复用的 Skill*")
st.markdown("---")

tab1, tab2 = st.tabs(["🆕 创建新 Skill", "📂 使用已有 Skill"])

# ===== 标签页1：创建新 Skill =====
with tab1:

    st.markdown("### 第一步：描述你的任务")

    col1, col2 = st.columns(2)
    with col1:
        task_desc = st.text_area("任务描述", placeholder="例如：写一篇小红书笔记", height=100)
    with col2:
        deliverable = st.text_area("交付要求", placeholder="例如：500字左右，包含标题、正文、标签", height=100)

    # 文件上传
    st.markdown("### 📎 上传参考文件（可选）")
    uploaded_files = st.file_uploader(
        "支持 txt、docx、xlsx、pptx、pdf、csv、json、md 等格式，可以上传多个文件",
        accept_multiple_files=True,
        type=["txt", "md", "json", "csv", "docx", "xlsx", "xls", "pptx", "pdf"]
    )

    if uploaded_files:
        all_text = ""
        for uf in uploaded_files:
            st.markdown(f"✅ 已上传：**{uf.name}**")
            file_text = read_uploaded_file(uf)
            all_text += f"\n\n=== 文件：{uf.name} ===\n{file_text}"
        st.session_state.uploaded_text = all_text
        with st.expander("📄 查看提取的文件内容"):
            st.text(all_text[:3000] + ("..." if len(all_text) > 3000 else ""))

    if st.button("🚀 生成 SOP", type="primary", use_container_width=True):
        if not task_desc.strip() or not deliverable.strip():
            st.error("请填写任务描述和交付要求")
        else:
            full_task = task_desc
            if st.session_state.uploaded_text:
                full_task += f"\n\n## 参考资料\n{st.session_state.uploaded_text}"
            with st.spinner("正在生成 SOP..."):
                try:
                    sop = call_generate_sop(full_task, deliverable)
                    st.session_state.sop = sop
                    st.session_state.sop_history = []
                    st.success("SOP 生成成功！")
                    st.rerun()
                except Exception as e:
                    st.error(f"生成失败：{e}")

    if st.session_state.sop is not None:

        st.markdown("---")
        st.markdown("### 第二步：审核和修改 SOP")
        display_sop(st.session_state.sop)

        st.markdown("---")
        feedback = st.text_input("修改意见", placeholder="例如：第三步太笼统了，请拆成更细的步骤")

        col_a, col_b = st.columns(2)
        with col_a:
            if st.button("✏️ 提交修改", use_container_width=True):
                if not feedback.strip():
                    st.warning("请输入修改意见")
                else:
                    with st.spinner("正在修改 SOP..."):
                        try:
                            st.session_state.sop_history.append(st.session_state.sop)
                            new_sop = call_refine_sop(st.session_state.sop, feedback)
                            st.session_state.sop = new_sop
                            st.success("修改成功！")
                            st.rerun()
                        except Exception as e:
                            st.session_state.sop_history.pop()
                            st.error(f"修改失败：{e}")
        with col_b:
            if st.button("↩️ 撤销修改", use_container_width=True):
                if not st.session_state.sop_history:
                    st.warning("已经是最初版本，无法撤销")
                else:
                    st.session_state.sop = st.session_state.sop_history.pop()
                    st.success("已撤销！")
                    st.rerun()

        st.markdown("---")
        st.markdown("### 第三步：确认并生成 Skill")

        if st.button("✅ 确认 SOP，生成 Skill", type="primary", use_container_width=True):
            with st.spinner("正在生成 Skill（约需30秒）..."):
                try:
                    skill = call_generate_skill(st.session_state.sop)
                    filename = skill["skill_name"].replace(" ", "_").replace("/", "_")
                    filepath = os.path.join(SKILLS_DIR, f"{filename}.json")
                    with open(filepath, "w", encoding="utf-8") as f:
                        json.dump(skill, f, ensure_ascii=False, indent=2)
                    st.session_state.skill = skill
                    st.session_state.chat_history = []
                    st.success(f"Skill 已生成并保存到 {filepath}")
                    st.rerun()
                except Exception as e:
                    st.error(f"生成失败：{e}")

    if st.session_state.skill is not None:

        st.markdown("---")
        st.markdown("### 第四步：复制 System Prompt 到其他平台")
        st.text_area(
            "System Prompt（复制到 ChatGPT / Coze / Dify 使用）",
            value=st.session_state.skill["system_prompt"],
            height=200,
            key="prompt_copy"
        )

        st.markdown("---")
        st.markdown("### 第五步：多轮对话试用 Skill")

        # 上传文件作为对话参考
        chat_files = st.file_uploader(
            "📎 上传文件作为输入（可选）",
            accept_multiple_files=True,
            type=["txt", "md", "json", "csv", "docx", "xlsx", "xls", "pptx", "pdf"],
            key="chat_files_tab1"
        )

        chat_file_text = ""
        if chat_files:
            for cf in chat_files:
                st.markdown(f"✅ 已上传：**{cf.name}**")
                chat_file_text += f"\n\n=== 文件：{cf.name} ===\n{read_uploaded_file(cf)}"

        # 选择输出格式
        output_format = st.selectbox(
            "📤 输出格式",
            ["纯文字（不生成文件）", "Word (.docx)", "Excel (.xlsx)", "PPT (.pptx)", "TXT (.txt)", "Markdown (.md)", "JSON (.json)"],
            key="output_format_tab1"
        )

        # 显示历史对话
        for msg in st.session_state.chat_history:
            if msg["role"] == "user":
                st.chat_message("user").markdown(msg["content"])
            else:
                st.chat_message("assistant").markdown(msg["content"])

        # 输入框
        user_msg = st.chat_input("输入你的内容（例如：主题是邪修过年）")

        if user_msg:
            full_msg = user_msg
            if chat_file_text:
                full_msg += f"\n\n## 用户上传的参考文件内容\n{chat_file_text}"

            st.session_state.chat_history.append({"role": "user", "content": full_msg})
            st.chat_message("user").markdown(user_msg)

            with st.chat_message("assistant"):
                with st.spinner("思考中..."):
                    try:
                        messages = [{"role": "system", "content": st.session_state.skill["system_prompt"]}]
                        messages.extend(st.session_state.chat_history)

                        response = client.chat.completions.create(
                            model="deepseek-chat",
                            messages=messages,
                            temperature=0.3
                        )
                        reply = response.choices[0].message.content
                        st.markdown(reply)
                        st.session_state.chat_history.append({"role": "assistant", "content": reply})

                        # 如果选了文件格式，生成下载按钮
                        if output_format != "纯文字（不生成文件）":
                            fmt_map = {
                                "Word (.docx)": "docx",
                                "Excel (.xlsx)": "xlsx",
                                "PPT (.pptx)": "pptx",
                                "TXT (.txt)": "txt",
                                "Markdown (.md)": "md",
                                "JSON (.json)": "json"
                            }
                            fmt = fmt_map.get(output_format, "txt")
                            file_data, file_name, mime_type = auto_generate_file(
                                reply, fmt, st.session_state.skill["skill_name"]
                            )
                            st.download_button(
                                label=f"📥 下载 {file_name}",
                                data=file_data,
                                file_name=file_name,
                                mime=mime_type
                            )

                    except Exception as e:
                        st.error(f"执行失败：{e}")

        if st.session_state.chat_history:
            if st.button("🗑️ 清空对话，重新开始", key="clear_tab1"):
                st.session_state.chat_history = []
                st.rerun()


# ===== 标签页2：使用已有 Skill =====
with tab2:

    st.markdown("### 加载之前保存的 Skill")

    skill_files = []
    if os.path.exists(SKILLS_DIR):
        skill_files = [f.replace(".json", "").replace("_", " ")
                       for f in sorted(os.listdir(SKILLS_DIR)) if f.endswith(".json")]

    if not skill_files:
        st.info("还没有保存过 Skill，请先在「创建新 Skill」标签页创建一个")
    else:
        selected = st.selectbox("选择 Skill", skill_files)

        if st.button("📥 加载 Skill", type="primary"):
            filename = selected.replace(" ", "_") + ".json"
            filepath = os.path.join(SKILLS_DIR, filename)
            try:
                with open(filepath, "r", encoding="utf-8") as f:
                    skill = json.load(f)
                st.session_state.skill = skill
                st.session_state.chat_history = []
                st.success(f"已加载：{selected}")
                st.rerun()
            except Exception as e:
                st.error(f"加载失败：{e}")

        if st.session_state.skill is not None:

            st.markdown("---")
            st.markdown(f"**当前 Skill：** {st.session_state.skill['skill_name']}")
            st.markdown(f"**描述：** {st.session_state.skill['description']}")

            st.text_area(
                "System Prompt（复制到其他平台使用）",
                value=st.session_state.skill["system_prompt"],
                height=200,
                key="prompt_copy_tab2"
            )

            st.markdown("---")
            st.markdown("### 多轮对话使用 Skill")

            # 上传文件
            chat_files2 = st.file_uploader(
                "📎 上传文件作为输入（可选）",
                accept_multiple_files=True,
                type=["txt", "md", "json", "csv", "docx", "xlsx", "xls", "pptx", "pdf"],
                key="chat_files_tab2"
            )

            chat_file_text2 = ""
            if chat_files2:
                for cf in chat_files2:
                    st.markdown(f"✅ 已上传：**{cf.name}**")
                    chat_file_text2 += f"\n\n=== 文件：{cf.name} ===\n{read_uploaded_file(cf)}"

            # 选择输出格式
            output_format2 = st.selectbox(
                "📤 输出格式",
                ["纯文字（不生成文件）", "Word (.docx)", "Excel (.xlsx)", "PPT (.pptx)", "TXT (.txt)", "Markdown (.md)", "JSON (.json)"],
                key="output_format_tab2"
            )

            # 显示历史对话
            for msg in st.session_state.chat_history:
                if msg["role"] == "user":
                    st.chat_message("user").markdown(msg["content"])
                else:
                    st.chat_message("assistant").markdown(msg["content"])

            user_msg2 = st.chat_input("输入你的需求...", key="chat_tab2")

            if user_msg2:
                full_msg2 = user_msg2
                if chat_file_text2:
                    full_msg2 += f"\n\n## 用户上传的参考文件内容\n{chat_file_text2}"

                st.session_state.chat_history.append({"role": "user", "content": full_msg2})
                st.chat_message("user").markdown(user_msg2)

                with st.chat_message("assistant"):
                    with st.spinner("思考中..."):
                        try:
                            messages = [{"role": "system", "content": st.session_state.skill["system_prompt"]}]
                            messages.extend(st.session_state.chat_history)

                            response = client.chat.completions.create(
                                model="deepseek-chat",
                                messages=messages,
                                temperature=0.3
                            )
                            reply = response.choices[0].message.content
                            st.markdown(reply)
                            st.session_state.chat_history.append({"role": "assistant", "content": reply})

                            if output_format2 != "纯文字（不生成文件）":
                                fmt_map = {
                                    "Word (.docx)": "docx",
                                    "Excel (.xlsx)": "xlsx",
                                    "PPT (.pptx)": "pptx",
                                    "TXT (.txt)": "txt",
                                    "Markdown (.md)": "md",
                                    "JSON (.json)": "json"
                                }
                                fmt = fmt_map.get(output_format2, "txt")
                                file_data, file_name, mime_type = auto_generate_file(
                                    reply, fmt, st.session_state.skill["skill_name"]
                                
                                )
                                st.download_button(
                                    label=f"📥 下载 {file_name}",
                                    data=file_data,
                                    file_name=file_name,
                                    mime=mime_type,
                                    key=f"dl_{file_name}"
                                )

                        except Exception as e:
                            st.error(f"执行失败：{e}")

            if st.session_state.chat_history:
                if st.button("🗑️ 清空对话，重新开始", key="clear_tab2"):
                    st.session_state.chat_history = []
                    st.rerun()

